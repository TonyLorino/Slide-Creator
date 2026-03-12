#!/usr/bin/env python3
"""
Extract all layout definitions from the Digital Realty PPTX template.
Outputs a JSON file with precise placeholder data for each of the 56 layouts.
"""

import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Emu

TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "public", "template", "PPT_Template.pptx",
)

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "lib", "templates", "layouts.json",
)


def emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 3)


def safe_color(font):
    try:
        if font.color and font.color.type is not None:
            return str(font.color.rgb)
    except Exception:
        pass
    try:
        if font.color and font.color.theme_color is not None:
            return f"theme:{font.color.theme_color}"
    except Exception:
        pass
    return None


def safe_fill_color(background):
    try:
        fill = background.fill
        if fill.type is not None:
            if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.type is not None:
                return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def is_placeholder(shape):
    try:
        return shape.placeholder_format is not None
    except (ValueError, AttributeError):
        return False


def extract_placeholder(shape):
    pf = shape.placeholder_format
    ph_type_map = {
        0: "title", 1: "body", 2: "body", 3: "title",
        4: "subtitle", 5: "body", 6: "body", 7: "body",
        8: "body", 9: "body", 10: "body", 11: "body",
        12: "slideNumber", 13: "body", 14: "body",
        15: "body", 16: "body", 17: "body", 18: "picture",
    }

    raw_type = pf.type
    raw_type_int = int(raw_type) if raw_type is not None else -1
    name = shape.name or ""
    name_lower = name.lower()
    idx = pf.idx

    # Name-based detection takes priority for accuracy
    if "picture" in name_lower:
        ph_type = "picture"
    elif "slide number" in name_lower:
        ph_type = "slideNumber"
    elif "subtitle" in name_lower:
        ph_type = "subtitle"
    elif "title" in name_lower and idx == 0:
        ph_type = "title"
    elif "title" in name_lower:
        ph_type = "title"
    elif raw_type_int == 18:
        ph_type = "picture"
    elif raw_type_int == 12:
        ph_type = "slideNumber"
    elif raw_type_int in (3, 15) or idx == 0:
        ph_type = "title"
    elif raw_type_int == 4:
        ph_type = "subtitle"
    else:
        ph_type = ph_type_map.get(raw_type_int, "body")

    default_text = ""
    font_size = None
    if shape.has_text_frame:
        parts = []
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
            for run in para.runs:
                if run.font.size and not font_size:
                    font_size = round(run.font.size.pt, 1)
        default_text = "\n".join(parts)

    return {
        "idx": pf.idx,
        "type": ph_type,
        "name": name,
        "x": emu_to_inches(shape.left),
        "y": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "defaultText": default_text if default_text else None,
        "fontSize": font_size,
    }


def extract_shape(shape):
    shape_type_map = {
        1: "auto", 2: "callout", 3: "canvas", 4: "chart",
        5: "freeform", 6: "group", 7: "embedded", 8: "line",
        9: "linked", 10: "linked_embedded", 11: "media",
        12: "placeholder", 13: "picture", 14: "script",
        16: "table", 17: "text",
    }
    raw_type = int(shape.shape_type) if shape.shape_type is not None else -1
    stype = shape_type_map.get(raw_type, "unknown")

    if "logo" in (shape.name or "").lower() or "graphic" in (shape.name or "").lower():
        stype = "logo"

    return {
        "name": shape.name or "",
        "type": stype,
        "x": emu_to_inches(shape.left),
        "y": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
    }


def categorize_layout(index: int, name: str, bg: str) -> str:
    nl = name.lower()
    if index <= 4 or "title" in nl and any(w in nl for w in ["1", "2", "3", "4"]) and "content" not in nl and "only" not in nl and "pattern" not in nl:
        if index <= 4:
            return "title"
    if "contents" in nl:
        return "contents"
    if "divider" in nl:
        return "divider"
    if "statement" in nl:
        return "statement"
    if "2 column" in nl or "two column" in nl:
        return "twoColumn"
    if "narrow" in nl:
        return "asymmetric"
    if "bold" in nl or "two line" in nl:
        return "boldTitle"
    if "copy" in nl and "image" in nl:
        return "copyImage"
    if "title only" in nl or "title with" in nl:
        return "titleOnly"
    if "blank" in nl:
        return "blank"
    if "close" in nl:
        return "close"
    if "slide content" in nl:
        return "content"
    return "content"


def extract_all():
    prs = Presentation(TEMPLATE_PATH)

    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    layouts = []

    for i, layout in enumerate(prs.slide_layouts):
        bg_color = safe_fill_color(layout.background) or "#FFFFFF"
        if bg_color and not bg_color.startswith("#"):
            bg_color = f"#{bg_color}"

        placeholders = []
        shapes = []

        for shape in layout.shapes:
            if is_placeholder(shape):
                placeholders.append(extract_placeholder(shape))
            else:
                shapes.append(extract_shape(shape))

        has_logo = any(
            s["type"] == "logo" or "graphic" in s["name"].lower()
            for s in shapes
        )

        logo_position = None
        if has_logo:
            logo_shape = next(
                (s for s in shapes if s["type"] == "logo" or "graphic" in s["name"].lower()),
                None,
            )
            if logo_shape:
                if logo_shape["y"] < 2:
                    logo_position = "titleTopLeft"
                elif logo_shape["y"] > 5:
                    logo_position = "closeBottomLeft"
                else:
                    logo_position = "contentFooter"

        has_slide_number = any(p["type"] == "slideNumber" for p in placeholders)

        has_corners = any("freeform" in s["name"].lower() for s in shapes)

        # Detect dark backgrounds from layout name or XML
        name = layout.name
        is_dark = "dark" in name.lower() or "black" in name.lower()

        # Try to detect specific known backgrounds
        if "colour" in name.lower() and "divider" in name.lower():
            bg_color = "#1F00FF"
        elif "light colour" in name.lower():
            bg_color = "#99F5FD"
        elif "statement 1" in name.lower() and "dark" not in name.lower():
            bg_color = "#9AB5B9"
        elif is_dark:
            bg_color = "#000000"
        elif bg_color == "#FFFFFF" or bg_color == "##FFFFFF":
            bg_color = "#FFFFFF"

        category = categorize_layout(i, name, bg_color)

        layouts.append({
            "index": i,
            "name": name,
            "category": category,
            "background": bg_color,
            "placeholders": placeholders,
            "shapes": shapes,
            "hasLogo": has_logo,
            "logoPosition": logo_position,
            "hasSlideNumber": has_slide_number,
            "hasCornerDecorations": has_corners,
        })

    result = {
        "slideWidth": slide_width,
        "slideHeight": slide_height,
        "layoutCount": len(layouts),
        "layouts": layouts,
    }

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    with open(OUTPUT_PATH, "w") as f:
        json.dump(result, f, indent=2)

    print(f"Extracted {len(layouts)} layouts to {OUTPUT_PATH}")
    print(f"Slide dimensions: {slide_width} x {slide_height} inches")

    for layout in layouts:
        ph_types = [p["type"] for p in layout["placeholders"]]
        print(f"  [{layout['index']:2d}] {layout['name']:<50} bg={layout['background']} ph={ph_types}")


if __name__ == "__main__":
    extract_all()
