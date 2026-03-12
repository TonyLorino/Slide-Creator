"""
Vercel Python Serverless Function for PPTX generation.
Receives deck JSON with the new simplified slide format, generates PowerPoint.
"""

import json
import io
import os
import re
from http.server import BaseHTTPRequestHandler
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(__file__))),
    "public",
    "template",
    "PPT_Template.pptx",
)

LAYOUT_NAME_TO_INDEX = {}
PH_INFO = {}


def _build_layout_map(prs):
    for layout in prs.slide_layouts:
        LAYOUT_NAME_TO_INDEX[layout.name] = layout

        phs = {}
        for shape in layout.shapes:
            try:
                pf = shape.placeholder_format
                if pf is not None:
                    name = (shape.name or "").lower()
                    idx = pf.idx
                    raw_type = int(pf.type) if pf.type is not None else -1

                    if "picture" in name:
                        ptype = "picture"
                    elif "slide number" in name:
                        ptype = "slideNumber"
                    elif "subtitle" in name:
                        ptype = "subtitle"
                    elif "title" in name:
                        ptype = "title"
                    elif raw_type == 18:
                        ptype = "picture"
                    elif raw_type == 12:
                        ptype = "slideNumber"
                    elif raw_type in (0, 3, 15):
                        ptype = "title"
                    elif raw_type == 4:
                        ptype = "subtitle"
                    else:
                        ptype = "body"

                    phs[idx] = {"type": ptype, "name": name}
            except (ValueError, AttributeError):
                pass

        PH_INFO[layout.name] = phs


def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
    return RGBColor(0, 0, 0)


_BULLET_RE = re.compile(r"^[\u2022\u2023\u25E6\u2043\u2219\-\*]\s*")


def _strip_bullets(text):
    lines = text.split("\n")
    return "\n".join(_BULLET_RE.sub("", line) for line in lines)


def _apply_text(ph, text, font_size=None, bold=None):
    if not text:
        return
    text = _strip_bullets(text)
    tf = ph.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = para_text
        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold


def _insert_image(ph, url):
    if not url:
        return
    try:
        if url.startswith("http"):
            import requests
            resp = requests.get(url, timeout=15)
            resp.raise_for_status()
            ph.insert_picture(io.BytesIO(resp.content))
        else:
            ph.insert_picture(url)
    except Exception:
        pass


def _add_table(slide, table_data):
    if not table_data or not table_data[0]:
        return
    nr, nc = len(table_data), len(table_data[0])
    left, top = Inches(0.6), Inches(2.0)
    width, height = Inches(12.0), Inches(4.5)
    shape = slide.shapes.add_table(nr, nc, left, top, width, height)
    table = shape.table
    for r in range(nr):
        for c in range(nc):
            cell = table.cell(r, c)
            cell.text = str(table_data[r][c]) if c < len(table_data[r]) else ""
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgb("01454F")
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
                        run.font.color.rgb = _hex_to_rgb("FFFFFF")
                        run.font.size = Pt(10)
            else:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)


def _add_chart(slide, chart_data):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if not chart_data:
        return
    labels = chart_data.get("labels", [])
    datasets = chart_data.get("datasets", [])
    if not labels or not datasets:
        return
    ct_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }
    cd = CategoryChartData()
    cd.categories = labels
    for ds in datasets:
        cd.add_series(ds.get("label", "Series"), ds.get("data", []))
    chart_type = ct_map.get(chart_data.get("chartType", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(chart_type, Inches(0.6), Inches(2.0), Inches(12.0), Inches(4.5), cd)


def generate_pptx(deck_data):
    prs = Presentation(TEMPLATE_PATH)
    _build_layout_map(prs)

    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    for sd in deck_data.get("slides", []):
        layout_key = sd.get("layoutKey", "Slide Content")
        layout = LAYOUT_NAME_TO_INDEX.get(layout_key, prs.slide_layouts[15])
        slide = prs.slides.add_slide(layout)

        ph_info = PH_INFO.get(layout_key, {})
        title_filled = False
        subtitle_filled = False
        body_count = 0
        filled_indices = set()

        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            info = ph_info.get(idx, {})
            ptype = info.get("type", "body")
            pname = info.get("name", "")

            if ptype == "slideNumber":
                filled_indices.add(idx)
                continue

            if ptype == "title" and not title_filled:
                title_filled = True
                if sd.get("title"):
                    _apply_text(ph, sd["title"], bold=True)
                    filled_indices.add(idx)

            elif ptype == "subtitle" and not subtitle_filled:
                subtitle_filled = True
                if sd.get("subtitle"):
                    _apply_text(ph, sd["subtitle"])
                    filled_indices.add(idx)

            elif ptype == "picture":
                url = sd.get("imageUrl", "")
                if url:
                    _insert_image(ph, url)
                    filled_indices.add(idx)

            elif ptype == "body":
                body_count += 1
                text = None
                if "presenter" in pname or "info" in pname:
                    text = sd.get("presenterInfo", "")
                elif body_count == 1:
                    text = sd.get("body", "")
                if text:
                    _apply_text(ph, text)
                    filled_indices.add(idx)

        # Remove unfilled placeholders so template defaults don't show
        for ph in list(slide.placeholders):
            idx = ph.placeholder_format.idx
            if idx not in filled_indices:
                sp = ph._element
                sp.getparent().remove(sp)

        if sd.get("tableData"):
            _add_table(slide, sd["tableData"])
        if sd.get("chartData"):
            _add_chart(slide, sd["chartData"])
        if sd.get("notes"):
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        content_length = int(self.headers.get("Content-Length", 0))
        body = self.rfile.read(content_length)

        try:
            deck_data = json.loads(body)
        except json.JSONDecodeError:
            self.send_response(400)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(json.dumps({"error": "Invalid JSON"}).encode())
            return

        try:
            pptx_bytes = generate_pptx(deck_data)
        except Exception as e:
            self.send_response(500)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(e)}).encode())
            return

        title = deck_data.get("title", "presentation").replace(" ", "_")
        filename = f"{title}.pptx"

        self.send_response(200)
        self.send_header(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
        self.send_header("Content-Length", str(len(pptx_bytes)))
        self.end_headers()
        self.wfile.write(pptx_bytes)
